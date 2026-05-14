#!/usr/bin/env python3
"""Convert reference materials (PDF, PPTX, DOCX) to markdown.

Usage:
    python -m scripts.convert_references reference/

Converts all supported files in the given directory to .md files
in the same directory. Skips files that already have a .md counterpart.

Dependencies (install via: uv sync):
    - pymupdf (fitz)  — PDF extraction
    - python-pptx     — PowerPoint extraction
    - python-docx     — Word document extraction
"""

import sys
from pathlib import Path


def convert_pdf(path: Path) -> str:
    """Extract text from PDF using pymupdf."""
    import fitz

    doc = fitz.open(path)
    parts = []
    for i, page in enumerate(doc, 1):
        text = page.get_text()
        if text.strip():
            parts.append(f"<!-- Page {i} -->\n\n{text.strip()}")
    doc.close()
    return "\n\n---\n\n".join(parts)


def convert_pptx(path: Path) -> str:
    """Extract text from PowerPoint slides."""
    from pptx import Presentation

    prs = Presentation(path)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    texts.append("| " + " | ".join(cells) + " |")
        if texts:
            slide_title = texts[0]
            body = "\n\n".join(texts[1:]) if len(texts) > 1 else ""
            parts.append(f"## Slide {i}: {slide_title}\n\n{body}")
    return "\n\n---\n\n".join(parts)


def convert_docx(path: Path) -> str:
    """Extract text from Word document."""
    from docx import Document

    doc = Document(path)
    parts = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style = para.style.name.lower() if para.style else ""
        if "heading 1" in style:
            parts.append(f"# {text}")
        elif "heading 2" in style:
            parts.append(f"## {text}")
        elif "heading 3" in style:
            parts.append(f"### {text}")
        else:
            parts.append(text)

    # Extract tables
    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append("| " + " | ".join(cells) + " |")
        if rows:
            # Add header separator after first row
            num_cols = len(table.rows[0].cells)
            header_sep = "| " + " | ".join(["---"] * num_cols) + " |"
            rows.insert(1, header_sep)
            parts.append("\n".join(rows))

    return "\n\n".join(parts)


CONVERTERS = {
    ".pdf": convert_pdf,
    ".pptx": convert_pptx,
    ".ppt": convert_pptx,  # python-pptx only supports .pptx; will error on old .ppt
    ".docx": convert_docx,
    ".doc": convert_docx,  # python-docx only supports .docx; will error on old .doc
}


def main():
    if len(sys.argv) < 2:
        print(f"Usage: {sys.argv[0]} <directory>")
        sys.exit(1)

    ref_dir = Path(sys.argv[1])
    if not ref_dir.is_dir():
        print(f"Error: {ref_dir} is not a directory")
        sys.exit(1)

    files = sorted(ref_dir.iterdir())
    converted = 0
    skipped = 0
    errors = []

    for f in files:
        if f.suffix.lower() not in CONVERTERS:
            continue

        out_path = f.with_suffix(".md")
        if out_path.exists():
            print(f"  skip (exists): {out_path.name}")
            skipped += 1
            continue

        print(f"  converting: {f.name} -> {out_path.name}")
        try:
            converter = CONVERTERS[f.suffix.lower()]
            md_content = converter(f)
            header = f"# {f.stem}\n\n> Converted from `{f.name}`\n\n"
            out_path.write_text(header + md_content, encoding="utf-8")
            converted += 1
        except Exception as e:
            errors.append((f.name, str(e)))
            print(f"  ERROR: {f.name}: {e}")

    print(f"\nDone: {converted} converted, {skipped} skipped, {len(errors)} errors")
    if errors:
        print("\nErrors:")
        for name, err in errors:
            print(f"  {name}: {err}")
        sys.exit(1)


if __name__ == "__main__":
    main()
